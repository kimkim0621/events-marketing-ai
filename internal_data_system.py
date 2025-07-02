#!/usr/bin/env python3
"""
社内データ統合システム
- 既存CSV/Excel読み込み
- PDF解析・情報抽出
- 知見データベース管理
- メディア属性管理
"""

import pandas as pd
import json
import sqlite3
import re
import os
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Any, Optional
import argparse

# PDF処理ライブラリ（オプション）
try:
    import PyPDF2
    import pdfplumber
    PDF_AVAILABLE = True
except ImportError:
    print("📋 PDF処理用: pip install PyPDF2 pdfplumber")
    PDF_AVAILABLE = False

# Claude API（オプション）
try:
    import anthropic
    CLAUDE_AVAILABLE = True
except ImportError:
    print("📋 Claude API用: pip install anthropic")
    CLAUDE_AVAILABLE = False

# PowerPoint処理ライブラリ（オプション）
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    print("📋 PowerPoint処理用: pip install python-pptx")
    PPTX_AVAILABLE = False

# Word文書処理ライブラリ（オプション）
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    print("📋 Word文書処理用: pip install python-docx")
    DOCX_AVAILABLE = False

class InternalDataSystem:
    """社内データ統合管理システム"""
    
    def __init__(self, db_path: str = None):
        if db_path is None:
            # クラウド対応のデータベースパス設定
            if "STREAMLIT_CLOUD" in os.environ or not os.path.exists("data"):
                self.db_path = "events_marketing.db"
            else:
                self.db_path = "data/events_marketing.db"
        else:
            self.db_path = db_path
        
        # Claude APIクライアントの初期化
        self.claude_client = None
        if CLAUDE_AVAILABLE:
            try:
                api_key = os.getenv('ANTHROPIC_API_KEY')
                if api_key:
                    self.claude_client = anthropic.Anthropic(api_key=api_key)
            except Exception as e:
                print(f"⚠️ Claude API初期化エラー: {e}")
        
        self.ensure_tables()
    
    def ensure_tables(self):
        """データベーステーブルの確保"""
        conn = sqlite3.connect(self.db_path)
        cursor = conn.cursor()
        
        # 知見データベース
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS internal_knowledge (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                category TEXT NOT NULL,
                title TEXT NOT NULL,
                content TEXT NOT NULL,
                conditions TEXT,
                impact_score REAL DEFAULT 1.0,
                confidence REAL DEFAULT 0.8,
                source TEXT,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        ''')
        
        # メディア詳細属性
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS media_detailed_attributes (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                media_name TEXT NOT NULL,
                attribute_category TEXT NOT NULL,
                attribute_name TEXT NOT NULL,
                attribute_value TEXT NOT NULL,
                data_source TEXT,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        ''')
        
        # 新規テーブル: メディア基本情報
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS media_basic_info (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                media_name TEXT NOT NULL UNIQUE,
                media_type TEXT,
                target_audience TEXT,
                description TEXT,
                website_url TEXT,
                contact_info TEXT,
                data_source TEXT,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        ''')
        
        conn.commit()
        conn.close()
    
    def import_existing_csv(self, file_path: str, data_type: str = "events") -> Dict:
        """既存CSVファイルのインポート（改善版）"""
        print(f"📊 既存データ読み込み: {file_path}")
        
        try:
            # CSV読み込みの改善（複数のエンコーディングを試行）
            df = None
            encodings = ['utf-8-sig', 'utf-8', 'shift_jis', 'cp932', 'latin1']
            
            for encoding in encodings:
                try:
                    df = pd.read_csv(file_path, encoding=encoding)
                    print(f"✅ エンコーディング {encoding} で読み込み成功")
                    break
                except UnicodeDecodeError:
                    continue
                except Exception as e:
                    print(f"⚠️ エンコーディング {encoding} で読み込み失敗: {e}")
                    continue
            
            if df is None:
                return {"success": False, "error": "ファイルを読み込めませんでした（エンコーディングエラー）"}
            
            print(f"📋 {len(df)}行 x {len(df.columns)}列を検出")
            print(f"🔍 列: {list(df.columns)}")
            
            # 空の列名を修正
            if any(col.startswith('Unnamed:') for col in df.columns):
                print("⚠️ ヘッダー行が検出されませんでした。最初の行をヘッダーとして使用します")
                # 最初のデータ行をヘッダーとして使用
                if len(df) > 0:
                    first_row = df.iloc[0]
                    new_columns = []
                    for i, val in enumerate(first_row):
                        if pd.notna(val) and str(val).strip():
                            new_columns.append(str(val).strip())
                        else:
                            new_columns.append(f"Column_{i+1}")
                    df.columns = new_columns
                    df = df.drop(df.index[0]).reset_index(drop=True)
                    print(f"🔧 新しい列名: {list(df.columns)}")
            
            # データタイプに基づく処理
            if data_type == "events":
                return self._process_event_csv(df, file_path)
            elif data_type == "media":
                return self._process_media_csv(df, file_path)
            elif data_type == "knowledge":
                return self._process_knowledge_csv(df, file_path)
            else:
                return {"success": False, "error": f"不明なデータタイプ: {data_type}"}
            
        except Exception as e:
            return {"success": False, "error": f"ファイル読み込みエラー: {str(e)}"}
    
    def _process_event_csv(self, df: pd.DataFrame, source: str) -> Dict:
        """イベントCSVの処理（改善版）"""
        conn = sqlite3.connect(self.db_path)
        cursor = conn.cursor()
        
        # 柔軟な列マッピング（更新版）
        mappings = {
            # イベント名
            'イベント名': 'event_name',
            'Event Name': 'event_name', 
            'event_name': 'event_name',
            'イベント': 'event_name',
            'Event': 'event_name',
            'Name': 'event_name',
            '名前': 'event_name',
            
            # カテゴリ・テーマ
            'カテゴリ': 'category',
            'Category': 'category',
            'category': 'category',
            '種類': 'category',
            'Type': 'category',
            'テーマ': 'theme',
            'Theme': 'theme',
            'theme': 'theme',
            'テーマ・内容': 'theme',
            
            # 参加者数
            '目標参加者数': 'target_attendees',
            '目標申込数': 'target_attendees',
            'Target': 'target_attendees',
            'target_attendees': 'target_attendees',
            '目標': 'target_attendees',
            
            '実際参加者数': 'actual_attendees',
            '実際申込数': 'actual_attendees',
            'Actual': 'actual_attendees',
            'actual_attendees': 'actual_attendees',
            '実績': 'actual_attendees',
            '結果': 'actual_attendees',
            
            # 予算・コスト
            '予算': 'budget',
            'Budget': 'budget',
            'budget': 'budget',
            'コスト': 'actual_cost',
            'Cost': 'actual_cost',
            'actual_cost': 'actual_cost',
            '実際コスト': 'actual_cost',
            '費用': 'actual_cost',
            
            # 日付
            '開催日': 'event_date',
            'Date': 'event_date',
            'event_date': 'event_date',
            '日付': 'event_date',
            
            # 施策
            '使用施策': 'campaigns',
            'Campaigns': 'campaigns',
            'campaigns': 'campaigns',
            '施策': 'campaigns',
        }
        
        df_mapped = df.rename(columns=mappings)
        imported = 0
        errors = []
        
        for index, row in df_mapped.iterrows():
            try:
                # 必須フィールドの処理
                event_name = str(row.get('event_name', f'Event_{imported+1}')).strip()
                if not event_name or event_name == 'nan':
                    event_name = f'インポートイベント_{imported+1}'
                
                # カテゴリの処理
                category = str(row.get('category', 'seminar')).strip()
                if not category or category == 'nan':
                    category = 'seminar'
                
                # テーマの処理（必須フィールド）
                theme = str(row.get('theme', 'その他')).strip()
                if not theme or theme == 'nan':
                    theme = 'その他'
                
                # 数値フィールドの処理
                try:
                    target = int(float(row.get('target_attendees', 0) or 0))
                except (ValueError, TypeError):
                    target = 0
                
                try:
                    actual = int(float(row.get('actual_attendees', 0) or 0))
                except (ValueError, TypeError):
                    actual = 0
                
                try:
                    budget = int(float(row.get('budget', 0) or 0))
                except (ValueError, TypeError):
                    budget = 0
                
                try:
                    cost = int(float(row.get('actual_cost', budget) or budget))
                except (ValueError, TypeError):
                    cost = budget
                
                # 日付の処理
                event_date = str(row.get('event_date', '2025-01-01')).strip()
                if not event_date or event_date == 'nan':
                    event_date = datetime.now().strftime('%Y-%m-%d')
                
                # 施策データの処理
                campaigns = row.get('campaigns', 'email_marketing')
                if pd.isna(campaigns) or campaigns == '':
                    campaigns = 'email_marketing'
                
                if isinstance(campaigns, str) and ',' in campaigns:
                    campaigns = [c.strip() for c in campaigns.split(',')]
                campaigns_json = json.dumps(campaigns if isinstance(campaigns, list) else [str(campaigns)])
                
                # パフォーマンス計算
                conversion_rate = (actual / target * 100) if target > 0 else 0
                cpa = (cost / actual) if actual > 0 else 0
                
                performance = json.dumps({
                    "conversion_rate": conversion_rate,
                    "cpa": cpa,
                    "cost_efficiency": budget / cost if cost > 0 else 1
                })
                
                # データベースに挿入（themeフィールドを含む）
                cursor.execute('''
                    INSERT INTO historical_events 
                    (event_name, category, theme, target_attendees, actual_attendees, 
                     budget, actual_cost, event_date, campaigns_used, performance_metrics)
                    VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
                ''', (
                    event_name, category, theme, target, actual, 
                    budget, cost, event_date, campaigns_json, performance
                ))
                
                imported += 1
                
            except Exception as e:
                error_msg = f"行{index+1}: {str(e)}"
                errors.append(error_msg)
                print(f"⚠️ {error_msg}")
                continue
        
        conn.commit()
        conn.close()
        
        result = {"success": True, "imported": imported}
        if errors:
            result["errors"] = errors
            result["error_count"] = len(errors)
        
        return result
    
    def _process_media_csv(self, df: pd.DataFrame, source: str) -> Dict:
        """メディアCSVの処理（改善版）"""
        conn = sqlite3.connect(self.db_path)
        cursor = conn.cursor()
        
        # 柔軟な列マッピング（更新版）
        mappings = {
            # メディア名
            'メディア名': 'media_name',
            'Media Name': 'media_name',
            'media_name': 'media_name',
            'Name': 'media_name',
            '名前': 'media_name',
            
            # メディアタイプ
            'メディアタイプ': 'media_type',
            'Media Type': 'media_type',
            'media_type': 'media_type',
            'Type': 'media_type',
            '種類': 'media_type',
            
            # 対象読者・オーディエンス
            '対象読者': 'target_audience',
            'Target Audience': 'target_audience',
            'target_audience': 'target_audience',
            'Audience': 'target_audience',
            'ターゲット': 'target_audience',
            
            # パフォーマンス指標
            'CTR': 'ctr',
            'ctr': 'ctr',
            'クリック率': 'ctr',
            'CVR': 'cvr',
            'cvr': 'cvr',
            'コンバージョン率': 'cvr',
            'CPA': 'cpa',
            'cpa': 'cpa',
            '獲得単価': 'cpa',
            'コスト': 'cpa',
            
            # リーチ・規模
            'リーチ': 'reach',
            'Reach': 'reach',
            'reach': 'reach',
            'ユーザー数': 'reach',
            '読者数': 'reach',
            
            # 説明
            '説明': 'description',
            'Description': 'description',
            'description': 'description',
            '概要': 'description',
        }
        
        df_mapped = df.rename(columns=mappings)
        imported = 0
        errors = []
        
        for index, row in df_mapped.iterrows():
            try:
                # 必須フィールド: メディア名
                media_name = str(row.get('media_name', '')).strip()
                if not media_name or media_name == 'nan':
                    errors.append(f"行{index+1}: メディア名が必須です")
                    continue
                
                # メディアタイプの処理
                media_type = str(row.get('media_type', 'その他')).strip()
                if not media_type or media_type == 'nan':
                    media_type = 'その他'
                
                # 対象読者の処理
                target_audience = str(row.get('target_audience', '')).strip()
                if not target_audience or target_audience == 'nan':
                    target_audience = ''
                
                # 数値フィールドの処理
                try:
                    ctr = float(row.get('ctr', 2.0) or 2.0)
                except (ValueError, TypeError):
                    ctr = 2.0
                
                try:
                    cvr = float(row.get('cvr', 5.0) or 5.0)
                except (ValueError, TypeError):
                    cvr = 5.0
                
                try:
                    cpa = float(row.get('cpa', 5000) or 5000)
                except (ValueError, TypeError):
                    cpa = 5000
                
                try:
                    reach = int(float(row.get('reach', 10000) or 10000))
                except (ValueError, TypeError):
                    reach = 10000
                
                # 説明の処理
                description = str(row.get('description', '')).strip()
                if not description or description == 'nan':
                    description = ''
                
                # メディア基本情報の保存
                try:
                    cursor.execute('''
                        INSERT OR REPLACE INTO media_basic_info
                        (media_name, media_type, target_audience, description, data_source)
                        VALUES (?, ?, ?, ?, ?)
                    ''', (media_name, media_type, target_audience, description, source))
                except Exception as e:
                    # media_basic_infoテーブルが存在しない場合は作成
                    cursor.execute('''
                        CREATE TABLE IF NOT EXISTS media_basic_info (
                            id INTEGER PRIMARY KEY AUTOINCREMENT,
                            media_name TEXT NOT NULL UNIQUE,
                            media_type TEXT,
                            target_audience TEXT,
                            description TEXT,
                            website_url TEXT,
                            contact_info TEXT,
                            data_source TEXT,
                            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                        )
                    ''')
                    cursor.execute('''
                        INSERT OR REPLACE INTO media_basic_info
                        (media_name, media_type, target_audience, description, data_source)
                        VALUES (?, ?, ?, ?, ?)
                    ''', (media_name, media_type, target_audience, description, source))
                
                # メディアパフォーマンス情報の保存
                try:
                    cursor.execute('''
                        INSERT OR REPLACE INTO media_performance 
                        (media_name, ctr, cvr, cpa, reach)
                        VALUES (?, ?, ?, ?, ?)
                    ''', (media_name, ctr, cvr, cpa, reach))
                except Exception as e:
                    # media_performanceテーブルが存在しない場合は作成（簡易版）
                    cursor.execute('''
                        CREATE TABLE IF NOT EXISTS media_performance (
                            id INTEGER PRIMARY KEY AUTOINCREMENT,
                            media_name TEXT NOT NULL UNIQUE,
                            ctr REAL DEFAULT 2.0,
                            cvr REAL DEFAULT 5.0,
                            cpa REAL DEFAULT 5000,
                            reach INTEGER DEFAULT 10000,
                            updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                        )
                    ''')
                    cursor.execute('''
                        INSERT OR REPLACE INTO media_performance 
                        (media_name, ctr, cvr, cpa, reach)
                        VALUES (?, ?, ?, ?, ?)
                    ''', (media_name, ctr, cvr, cpa, reach))
                
                imported += 1
                
            except Exception as e:
                error_msg = f"行{index+1}: {str(e)}"
                errors.append(error_msg)
                print(f"⚠️ {error_msg}")
                continue
        
        conn.commit()
        conn.close()
        
        result = {"success": True, "imported": imported}
        if errors:
            result["errors"] = errors
            result["error_count"] = len(errors)
        
        return result
    
    def _process_knowledge_csv(self, df: pd.DataFrame, source: str) -> Dict:
        """知見CSVの処理"""
        conn = sqlite3.connect(self.db_path)
        cursor = conn.cursor()
        
        imported = 0
        
        for _, row in df.iterrows():
            try:
                category = str(row.get('category', row.get('カテゴリ', 'general')))
                title = str(row.get('title', row.get('タイトル', f'知見_{imported+1}')))
                content = str(row.get('content', row.get('内容', '')))
                
                if not content:
                    continue
                
                impact = float(row.get('impact_score', row.get('影響度', 1.0)) or 1.0)
                confidence = float(row.get('confidence', row.get('信頼度', 0.8)) or 0.8)
                
                cursor.execute('''
                    INSERT INTO internal_knowledge
                    (category, title, content, impact_score, confidence, source)
                    VALUES (?, ?, ?, ?, ?, ?)
                ''', (category, title, content, impact, confidence, source))
                
                imported += 1
                
            except Exception as e:
                print(f"⚠️ 知見行{imported+1}エラー: {e}")
                continue
        
        conn.commit()
        conn.close()
        
        return {"success": True, "imported": imported}
    
    def extract_pdf_insights(self, file_path: str) -> Dict:
        """PDFから知見・属性を抽出（改善版）"""
        if not PDF_AVAILABLE:
            return {"success": False, "error": "PDF処理ライブラリが必要"}
        
        print(f"📄 PDF解析: {file_path}")
        
        try:
            # PDFテキスト抽出
            with pdfplumber.open(file_path) as pdf:
                text = "\n".join([page.extract_text() or "" for page in pdf.pages])
            
            if not text.strip():
                return {"success": False, "error": "PDFからテキストを抽出できませんでした"}
            
            # Claude APIを使用した高精度分析
            if self.claude_client:
                analysis_result = self._analyze_pdf_with_claude(text, file_path)
            else:
                # フォールバック: 従来の方法
                analysis_result = self._analyze_pdf_fallback(text, file_path)
            
            return analysis_result
            
        except Exception as e:
            return {"success": False, "error": str(e)}
    
    def extract_pptx_insights(self, file_path: str) -> Dict:
        """PowerPointから知見・属性を抽出"""
        if not PPTX_AVAILABLE:
            return {"success": False, "error": "PowerPoint処理ライブラリが必要"}
        
        print(f"📊 PowerPoint解析: {file_path}")
        
        try:
            # PowerPointからテキスト抽出
            prs = Presentation(file_path)
            text_content = []
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = f"=== スライド {slide_num + 1} ===\n"
                
                # テキストボックスからテキスト抽出
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"
                    
                    # 表からテキスト抽出
                    if shape.shape_type == 19:  # Table
                        try:
                            table = shape.table
                            for row in table.rows:
                                row_text = []
                                for cell in row.cells:
                                    if cell.text.strip():
                                        row_text.append(cell.text.strip())
                                if row_text:
                                    slide_text += " | ".join(row_text) + "\n"
                        except:
                            pass
                
                text_content.append(slide_text)
            
            full_text = "\n".join(text_content)
            
            if not full_text.strip():
                return {"success": False, "error": "PowerPointからテキストを抽出できませんでした"}
            
            # Claude APIを使用した分析
            if self.claude_client:
                analysis_result = self._analyze_document_with_claude(full_text, file_path, "PowerPoint")
            else:
                # フォールバック: 従来の方法
                analysis_result = self._analyze_text_fallback(full_text, file_path, "PowerPoint")
            
            return analysis_result
            
        except Exception as e:
            return {"success": False, "error": str(e)}
    
    def extract_docx_insights(self, file_path: str) -> Dict:
        """Word文書から知見・属性を抽出"""
        if not DOCX_AVAILABLE:
            return {"success": False, "error": "Word文書処理ライブラリが必要"}
        
        print(f"📄 Word文書解析: {file_path}")
        
        try:
            # Word文書からテキスト抽出
            doc = Document(file_path)
            text_content = []
            
            # 段落からテキスト抽出
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    # 見出しスタイルの検出
                    if paragraph.style.name.startswith('Heading'):
                        text_content.append(f"\n## {paragraph.text}\n")
                    else:
                        text_content.append(paragraph.text)
            
            # 表からテキスト抽出
            for table in doc.tables:
                text_content.append("\n=== 表 ===")
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_content.append(" | ".join(row_text))
            
            full_text = "\n".join(text_content)
            
            if not full_text.strip():
                return {"success": False, "error": "Word文書からテキストを抽出できませんでした"}
            
            # Claude APIを使用した分析
            if self.claude_client:
                analysis_result = self._analyze_document_with_claude(full_text, file_path, "Word")
            else:
                # フォールバック: 従来の方法
                analysis_result = self._analyze_text_fallback(full_text, file_path, "Word")
            
            return analysis_result
            
        except Exception as e:
            return {"success": False, "error": str(e)}
    
    def _analyze_pdf_with_claude(self, text: str, source: str) -> Dict:
        """Claude APIを使用したPDF分析"""
        try:
            # Claude用のプロンプト設計
            prompt = f"""
以下のPDFテキストを分析し、内容を適切にカテゴリ分けして構造化データとして抽出してください。

PDFテキスト:
{text[:8000]}  # トークン数制限のため最初の8000文字まで

以下の形式でJSONを返してください:

{{
    "content_type": "media_info" または "knowledge_base" または "mixed",
    "confidence": 0.0-1.0,
    "media_information": [
        {{
            "media_name": "メディア名",
            "media_type": "技術系メディア/ビジネス系メディア/SNS/その他",
            "target_audience": "対象読者",
            "attributes": [
                {{
                    "category": "performance/audience/cost/general",
                    "name": "属性名",
                    "value": "属性値"
                }}
            ],
            "description": "メディアの説明"
        }}
    ],
    "knowledge_insights": [
        {{
            "category": "campaign/media/audience/budget/timing/general",
            "title": "知見のタイトル",
            "content": "知見の詳細内容",
            "impact_score": 0.0-1.0,
            "confidence": 0.0-1.0,
            "applicable_conditions": ["条件1", "条件2"]
        }}
    ]
}}

注意:
- メディア情報と知見を明確に区別してください
- メディア情報: 媒体の基本情報、パフォーマンス指標、読者層など
- 知見: ノウハウ、ベストプラクティス、注意点、推奨事項など
- 抽出できない場合は空の配列を返してください
- confidenceは抽出精度の自信度を表してください
"""
            
            message = self.claude_client.messages.create(
                model="claude-3-5-sonnet-20241022",
                max_tokens=4000,
                temperature=0.1,
                messages=[
                    {
                        "role": "user",
                        "content": prompt
                    }
                ]
            )
            
            # レスポンスの解析
            response_text = message.content[0].text
            
            # JSONの抽出（マークダウンコードブロック対応）
            json_match = re.search(r'```(?:json)?\s*(\{.*?\})\s*```', response_text, re.DOTALL)
            if json_match:
                json_text = json_match.group(1)
            else:
                json_text = response_text
            
            analysis = json.loads(json_text)
            
            # データベースへの保存
            media_extracted = self._save_media_info(analysis.get('media_information', []), source)
            insights_extracted = self._save_insights(analysis.get('knowledge_insights', []), source)
            
            return {
                "success": True,
                "content_type": analysis.get('content_type', 'unknown'),
                "confidence": analysis.get('confidence', 0.0),
                "media_extracted": media_extracted,
                "insights_extracted": insights_extracted,
                "analysis_method": "claude_api"
            }
            
        except Exception as e:
            print(f"⚠️ Claude分析エラー: {e}")
            # フォールバックに切り替え
            return self._analyze_pdf_fallback(text, source)
    
    def _analyze_pdf_fallback(self, text: str, source: str) -> Dict:
        """従来の正規表現ベース分析（フォールバック）"""
        try:
            # 既存の処理を改善
            media_info = self._extract_media_from_pdf(text, source)
            insights = self._extract_insights_from_pdf(text, source)
            
            return {
                "success": True,
                "content_type": "mixed",
                "confidence": 0.5,
                "media_extracted": media_info,
                "insights_extracted": insights,
                "analysis_method": "regex_fallback"
            }
            
        except Exception as e:
            return {"success": False, "error": str(e)}
    
    def _save_media_info(self, media_info_list, source):
        """メディア情報をデータベースに保存"""
        if not media_info_list:
            return 0
        
        conn = sqlite3.connect(self.db_path)
        cursor = conn.cursor()
        
        saved_count = 0
        
        for media_info in media_info_list:
            try:
                media_name = media_info.get('media_name', '')
                media_type = media_info.get('media_type', '')
                target_audience = media_info.get('target_audience', '')
                description = media_info.get('description', '')
                
                if not media_name:
                    continue
                
                # メディア基本情報を保存
                cursor.execute('''
                    INSERT OR REPLACE INTO media_basic_info
                    (media_name, media_type, target_audience, description, data_source)
                    VALUES (?, ?, ?, ?, ?)
                ''', (media_name, media_type, target_audience, description, source))
                
                # 属性情報を保存
                attributes = media_info.get('attributes', [])
                for attr in attributes:
                    attr_category = attr.get('category', 'general')
                    attr_name = attr.get('name', '')
                    attr_value = attr.get('value', '')
                    
                    if attr_name and attr_value:
                        cursor.execute('''
                            INSERT OR REPLACE INTO media_detailed_attributes
                            (media_name, attribute_category, attribute_name, attribute_value, data_source)
                            VALUES (?, ?, ?, ?, ?)
                        ''', (media_name, attr_category, attr_name, attr_value, source))
                
                saved_count += 1
                
            except Exception as e:
                print(f"⚠️ メディア情報保存エラー: {e}")
                continue
        
        conn.commit()
        conn.close()
        
        return saved_count
    
    def _save_insights(self, insights_list, source):
        """知見をデータベースに保存"""
        if not insights_list:
            return 0
        
        conn = sqlite3.connect(self.db_path)
        cursor = conn.cursor()
        
        saved_count = 0
        
        for insight in insights_list:
            try:
                category = insight.get('category', 'general')
                title = insight.get('title', '')
                content = insight.get('content', '')
                impact_score = insight.get('impact_score', 0.7)
                confidence = insight.get('confidence', 0.8)
                conditions = insight.get('applicable_conditions', [])
                
                if not title or not content:
                    continue
                
                # 条件をJSON形式で保存
                conditions_json = json.dumps({"general": conditions}) if conditions else None
                
                cursor.execute('''
                    INSERT INTO internal_knowledge
                    (category, title, content, impact_score, confidence, conditions, source)
                    VALUES (?, ?, ?, ?, ?, ?, ?)
                ''', (category, title, content, impact_score, confidence, conditions_json, source))
                
                saved_count += 1
                
            except Exception as e:
                print(f"⚠️ 知見保存エラー: {e}")
                continue
        
        conn.commit()
        conn.close()
        
        return saved_count
    
    def _extract_media_from_pdf(self, text: str, source: str) -> int:
        """PDFからメディア属性を抽出"""
        conn = sqlite3.connect(self.db_path)
        cursor = conn.cursor()
        
        extracted = 0
        
        # メディア名の検出パターン
        media_patterns = [
            r'([A-Za-z\s]+(?:Tech|IT|Engineer|Developer|Code|Programming))\s*(?:媒体|メディア)',
            r'媒体[:：]\s*([^\n\r]+)',
        ]
        
        # 属性抽出パターン
        attribute_patterns = {
            'target_jobs': r'(?:対象職種|職種)[:：]\s*([^\n\r]+)',
            'target_industries': r'(?:対象業界|業界)[:：]\s*([^\n\r]+)', 
            'ctr_rate': r'CTR[:：]\s*([0-9.]+)%?',
            'conversion_rate': r'CV[R率][:：]\s*([0-9.]+)%?',
            'cost_per_click': r'CP[CA][:：]\s*([0-9,]+)円?',
            'audience_size': r'(?:リーチ|読者数)[:：]\s*([0-9,]+)',
        }
        
        # メディア検出
        media_names = set()
        for pattern in media_patterns:
            matches = re.findall(pattern, text, re.IGNORECASE | re.MULTILINE)
            media_names.update([m.strip() for m in matches if len(m.strip()) > 2])
        
        # 各メディアの属性抽出
        for media_name in media_names:
            # メディア周辺のテキストを検索
            media_context = self._get_media_context(text, media_name)
            
            for attr_name, pattern in attribute_patterns.items():
                matches = re.findall(pattern, media_context, re.IGNORECASE)
                for match in matches:
                    cursor.execute('''
                        INSERT INTO media_detailed_attributes
                        (media_name, attribute_category, attribute_name, attribute_value, data_source)
                        VALUES (?, ?, ?, ?, ?)
                    ''', (media_name, self._categorize_attribute(attr_name), attr_name, match.strip(), source))
                    extracted += 1
        
        conn.commit()
        conn.close()
        
        return extracted
    
    def _extract_insights_from_pdf(self, text: str, source: str) -> int:
        """PDFから知見を抽出"""
        conn = sqlite3.connect(self.db_path)
        cursor = conn.cursor()
        
        extracted = 0
        
        # 知見パターンの検出
        insight_patterns = [
            r'(?:知見|ノウハウ|ベストプラクティス)[:：]\s*([^\n\r]+)',
            r'(?:効果的|有効)(?:な|である)\s*([^\n\r]+)',
            r'(?:推奨|おすすめ)[:：]\s*([^\n\r]+)',
            r'(?:注意|気をつける)べき(?:点|こと)[:：]\s*([^\n\r]+)',
        ]
        
        # カテゴリ判定パターン
        category_patterns = {
            'campaign': ['キャンペーン', '施策', 'マーケティング', '集客'],
            'media': ['メディア', '媒体', '広告', 'SNS'],
            'audience': ['ターゲット', 'オーディエンス', 'ユーザー', '参加者'],
            'budget': ['予算', 'コスト', '費用', '価格'],
            'timing': ['タイミング', '時期', 'スケジュール', '配信']
        }
        
        # 知見の抽出
        for pattern in insight_patterns:
            matches = re.findall(pattern, text, re.IGNORECASE | re.MULTILINE)
            for match in matches:
                insight_text = match.strip()
                if len(insight_text) > 10:  # 短すぎる内容を除外
                    
                    # カテゴリの推定
                    category = 'general'
                    for cat, keywords in category_patterns.items():
                        if any(keyword in insight_text for keyword in keywords):
                            category = cat
                            break
                    
                    # 重複チェック
                    cursor.execute('''
                        SELECT COUNT(*) FROM internal_knowledge 
                        WHERE content = ? AND source = ?
                    ''', (insight_text, source))
                    
                    if cursor.fetchone()[0] == 0:  # 重複なし
                        cursor.execute('''
                            INSERT INTO internal_knowledge
                            (category, title, content, impact_score, confidence, source)
                            VALUES (?, ?, ?, ?, ?, ?)
                        ''', (
                            category, 
                            f"PDF抽出知見_{extracted+1}",
                            insight_text,
                            0.7,  # PDF抽出は中程度の影響度
                            0.6,  # PDF抽出は中程度の信頼度
                            source
                        ))
                        extracted += 1
        
        conn.commit()
        conn.close()
        
        return extracted
    
    def _get_media_context(self, text: str, media_name: str, context_size: int = 500) -> str:
        """メディア名周辺のコンテキストを取得"""
        media_pos = text.lower().find(media_name.lower())
        if media_pos == -1:
            return text[:1000]  # 見つからない場合は先頭を返す
        
        start = max(0, media_pos - context_size)
        end = min(len(text), media_pos + len(media_name) + context_size)
        return text[start:end]
    
    def _categorize_attribute(self, attr_name: str) -> str:
        """属性のカテゴリ分類"""
        if 'job' in attr_name or 'industry' in attr_name:
            return 'audience'
        elif 'cost' in attr_name or 'ctr' in attr_name or 'conversion' in attr_name:
            return 'performance'
        elif 'size' in attr_name or 'reach' in attr_name:
            return 'scale'
        else:
            return 'general'
    
    def add_manual_knowledge(self, category: str, title: str, content: str, 
                           conditions: Dict = None, impact: float = 1.0) -> int:
        """手動での知見追加"""
        conn = sqlite3.connect(self.db_path)
        cursor = conn.cursor()
        
        cursor.execute('''
            INSERT INTO internal_knowledge
            (category, title, content, conditions, impact_score, source)
            VALUES (?, ?, ?, ?, ?, ?)
        ''', (
            category, title, content, 
            json.dumps(conditions) if conditions else None,
            impact, 'manual'
        ))
        
        knowledge_id = cursor.lastrowid
        conn.commit()
        conn.close()
        
        print(f"✅ 知見追加: {title} (ID: {knowledge_id})")
        return knowledge_id
    
    def get_applicable_knowledge(self, event_conditions: Dict) -> List[Dict]:
        """イベント条件に適用可能な知見を取得"""
        conn = sqlite3.connect(self.db_path)
        cursor = conn.cursor()
        
        cursor.execute('''
            SELECT * FROM internal_knowledge 
            ORDER BY impact_score DESC, confidence DESC
        ''')
        
        all_knowledge = cursor.fetchall()
        columns = [desc[0] for desc in cursor.description]
        
        applicable = []
        for row in all_knowledge:
            knowledge = dict(zip(columns, row))
            
            # 条件マッチング
            if knowledge['conditions']:
                try:
                    conditions = json.loads(knowledge['conditions'])
                    if self._matches_event_conditions(event_conditions, conditions):
                        applicable.append(knowledge)
                except:
                    continue
            else:
                # 汎用知見
                applicable.append(knowledge)
        
        conn.close()
        return applicable
    
    def _matches_event_conditions(self, event_cond: Dict, stored_cond: Dict) -> bool:
        """条件マッチング判定"""
        for key, value in stored_cond.items():
            if key not in event_cond:
                continue
            
            event_val = event_cond[key]
            if isinstance(value, list):
                if not any(v in str(event_val) for v in value):
                    return False
            elif str(value).lower() not in str(event_val).lower():
                return False
        return True
    
    def show_data_overview(self):
        """データ概要の表示"""
        conn = sqlite3.connect(self.db_path)
        cursor = conn.cursor()
        
        # 基本カウント
        tables = ['historical_events', 'media_performance', 'media_detailed_attributes', 'internal_knowledge']
        
        print("\n📊 社内データ概要")
        for table in tables:
            try:
                cursor.execute(f"SELECT COUNT(*) FROM {table}")
                count = cursor.fetchone()[0]
                print(f"  {table}: {count}件")
            except:
                print(f"  {table}: テーブルなし")
        
        # 知見カテゴリ別統計
        cursor.execute('''
            SELECT category, COUNT(*) FROM internal_knowledge 
            GROUP BY category ORDER BY COUNT(*) DESC
        ''')
        knowledge_stats = cursor.fetchall()
        
        if knowledge_stats:
            print("\n🧠 知見カテゴリ別")
            for category, count in knowledge_stats:
                print(f"  {category}: {count}件")
        
        conn.close()

def main():
    parser = argparse.ArgumentParser(description='社内データ統合システム')
    parser.add_argument('--import-csv', type=str, help='CSVファイルのインポート')
    parser.add_argument('--import-pdf', type=str, help='PDFファイルの解析')
    parser.add_argument('--data-type', choices=['events', 'media', 'knowledge'], 
                       default='events', help='データタイプ')
    parser.add_argument('--overview', action='store_true', help='データ概要表示')
    
    args = parser.parse_args()
    
    system = InternalDataSystem()
    
    if args.overview:
        system.show_data_overview()
        return
    
    if args.import_csv:
        result = system.import_existing_csv(args.import_csv, args.data_type)
        print(f"📊 CSVインポート結果: {result}")
    
    if args.import_pdf:
        result = system.extract_pdf_insights(args.import_pdf)
        print(f"📄 PDF解析結果: {result}")
    
    system.show_data_overview()

if __name__ == "__main__":
    main() 